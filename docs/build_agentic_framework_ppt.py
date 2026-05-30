"""
Generate a professional, light-themed PowerPoint deck from the
"Agentic framework - flow, plan, and patterns" canvas.

Run with the project venv:
    .venv\\Scripts\\python.exe docs\\build_agentic_framework_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Palette (light / professional, derived from the source canvas)
# ----------------------------------------------------------------------------
BG          = RGBColor(0xFB, 0xFA, 0xF7)   # warm off-white page
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x2C, 0x2C, 0x2A)   # primary text
INK2        = RGBColor(0x5F, 0x5E, 0x5A)   # secondary text
MUTED       = RGBColor(0x8A, 0x89, 0x80)

# Zone colours: (fill, border, deep-text)
INDIGO   = (RGBColor(0xEE, 0xED, 0xFE), RGBColor(0x53, 0x4A, 0xB7), RGBColor(0x3C, 0x34, 0x89))
TERRA    = (RGBColor(0xFA, 0xEC, 0xE7), RGBColor(0x99, 0x3C, 0x1D), RGBColor(0x71, 0x2B, 0x13))
AMBER    = (RGBColor(0xFA, 0xEE, 0xDA), RGBColor(0x85, 0x4F, 0x0B), RGBColor(0x63, 0x38, 0x06))
TEAL     = (RGBColor(0xE1, 0xF5, 0xEE), RGBColor(0x0F, 0x6E, 0x56), RGBColor(0x08, 0x50, 0x41))
NEUTRAL  = (RGBColor(0xF1, 0xEF, 0xE8), RGBColor(0x5F, 0x5E, 0x5A), RGBColor(0x2C, 0x2C, 0x2A))
GREEN    = (RGBColor(0xEA, 0xF3, 0xDE), RGBColor(0x3B, 0x6D, 0x11), RGBColor(0x27, 0x50, 0x0A))

ACCENT   = RGBColor(0x53, 0x4A, 0xB7)
ARROW    = RGBColor(0x6B, 0x6A, 0x64)

FONT = "Segoe UI"

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def _set_text(tf, lines, size, color, bold=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, space_after=2, line_spacing=1.0):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        # support (text, color, bold, size) tuples per-run
        if isinstance(ln, tuple):
            text, rc, rb, rs = (list(ln) + [color, bold, size])[:4]
        else:
            text, rc, rb, rs = ln, color, bold, size
        r = p.add_run(); r.text = text
        f = r.font
        f.name, f.size, f.bold = FONT, Pt(rs), rb
        f.color.rgb = rc


def textbox(s, x, y, w, h, lines, size=12, color=INK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=2,
            line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(tb.text_frame, lines, size, color, bold, align, anchor,
              space_after, line_spacing)
    return tb


def box(s, x, y, w, h, title=None, body=None, theme=NEUTRAL,
        title_size=14, body_size=11, rounded=True, align=PP_ALIGN.CENTER):
    fill, line, deep = theme
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(1.1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = []
    if title:
        lines.append((title, deep, True, title_size))
    if body:
        if isinstance(body, str):
            body = [body]
        for b in body:
            lines.append((b, line, False, body_size))
    _set_text(tf, lines, body_size, deep, False, align, MSO_ANCHOR.MIDDLE,
              space_after=2, line_spacing=1.0)
    return shp


def header(s, title, kicker=None):
    # accent side bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.5),
                             Inches(0.09), Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s, 0.78, 0.46, 11.8, 0.8, title, size=26, color=INK, bold=True)
    if kicker:
        textbox(s, 0.80, 1.12, 11.8, 0.4, kicker, size=13, color=INK2)
    # thin divider
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.6),
                                Inches(kicker and 1.55 or 1.32),
                                Inches(12.73), Inches(kicker and 1.55 or 1.32))
    ln.line.color.rgb = RGBColor(0xDD, 0xDB, 0xD2); ln.line.width = Pt(0.75)


def arrow(s, x1, y1, x2, y2, color=ARROW, width=1.6, dash=False):
    cxn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = color; cxn.line.width = Pt(width)
    cxn.shadow.inherit = False
    le = cxn.line._get_or_add_ln()
    tail = le.makeelement(qn('a:tailEnd'),
                          {'type': 'triangle', 'w': 'med', 'len': 'med'})
    le.append(tail)
    if dash:
        d = le.makeelement(qn('a:prstDash'), {'val': 'dash'})
        le.insert(0, d)
    return cxn


def chip(s, x, y, text, theme=INDIGO, w=None, h=0.34, size=10.5):
    fill, line, deep = theme
    if w is None:
        w = 0.22 + 0.085 * len(text)
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    _set_text(tf, text, size, deep, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return shp


def footer(s, n):
    textbox(s, 0.6, 7.05, 9, 0.3,
            "Agentic Framework  ·  Flow · Plan · Patterns", size=9, color=MUTED)
    textbox(s, 12.0, 7.05, 0.9, 0.3, str(n), size=9, color=MUTED,
            align=PP_ALIGN.RIGHT)


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.35), SW, Inches(2.7))
band.fill.solid(); band.fill.fore_color.rgb = WHITE
band.line.fill.background(); band.shadow.inherit = False
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.62),
                         Inches(0.12), Inches(2.15))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background(); bar.shadow.inherit = False
textbox(s, 1.35, 2.55, 11, 1.4,
        "Agentic Framework", size=46, color=INK, bold=True)
textbox(s, 1.38, 3.75, 11, 0.7,
        "Flow · Plan · Patterns — a reference architecture", size=22, color=ACCENT)
textbox(s, 1.4, 4.45, 11, 0.5,
        "Planner-driven DAG execution · shared blackboard · A2A transport · "
        "memory & observability", size=13.5, color=INK2)
# zone legend chips
chip(s, 1.4, 5.5, "Planning",  INDIGO, w=1.4)
chip(s, 2.95, 5.5, "Agents",   TERRA,  w=1.2)
chip(s, 4.3, 5.5, "Tools / Registry", AMBER, w=1.95)
chip(s, 6.4, 5.5, "Memory & State",  TEAL,  w=1.85)
chip(s, 8.4, 5.5, "Observability",   GREEN, w=1.7)
footer(s, 1)

# ============================================================================
# SLIDE 2 — Agenda
# ============================================================================
s = slide()
header(s, "What this deck covers", "Eight zones, one request-to-response lifecycle")
items = [
    ("01", "Request → response flow", "The end-to-end pipeline from prompt to answer", INDIGO),
    ("02", "How the Planner builds the plan", "Intent split, registry match, DAG, SLA budgets", INDIGO),
    ("03", "Plan → execution waves", "Dependencies become parallel waves", TERRA),
    ("04", "A2A transport", "JSON-RPC vs Kafka, derived not configured", AMBER),
    ("05", "Memory & state", "Scratchpad, Redis, Postgres, Milvus", TEAL),
    ("06", "Commit & write discipline", "Durable facts only after the gate passes", TEAL),
    ("07", "Observability & versioning", "Spans, traces, registry contract", GREEN),
    ("08", "Partial streaming & backpressure", "Progressive delivery, consumer-lag control", AMBER),
]
x0, y0, cw, ch, gx, gy = 0.75, 1.85, 5.95, 1.12, 0.35, 0.18
for i, (num, t, d, theme) in enumerate(items):
    col, row = i % 2, i // 2
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    b = box(s, x, y, cw, ch, theme=theme, rounded=True)
    b.text_frame.clear()
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, [(num + "   " + t, theme[2], True, 15),
                   (d, theme[1], False, 11.5)],
              12, theme[2], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
              space_after=3)
    tf.margin_left = Inches(0.2)
footer(s, 2)

# ============================================================================
# SLIDE 3 — Request -> response flow (the pipeline)
# ============================================================================
s = slide()
header(s, "Request → response flow",
       "One natural prompt drives the whole lifecycle — no agent calls another directly")
steps = [
    ("User request", "One natural prompt with several intents", NEUTRAL),
    ("Planner", "Split intents → match agents → wire DAG → assign SLA budgets", INDIGO),
    ("Orchestrator", "Runs the DAG in dependency waves · tracks agent health", INDIGO),
    ("Selected agents execute", "Wave 1 → Wave 2 (parallel) → Wave 3", TERRA),
    ("Tool / MCP gateway", "Outage DB, OASIS API, vector search · enforces bounds · emits spans", AMBER),
    ("Shared state / blackboard", "Every agent writes here · error flags visible to the gate", TEAL),
    ("Completion gate", "All done? Goal met? Bounds ok? → else re-plan (iter ≤ max)", INDIGO),
    ("Synthesizer", "Reads state, composes one answer · marks [PARTIAL] on errors", INDIGO),
    ("Final response", "A single, ordered answer for the user", NEUTRAL),
]
n = len(steps)
x, w = 0.9, 7.1
y = 1.78
h = 0.50
gap = 0.115
for i, (t, d, theme) in enumerate(steps):
    box(s, x, y, w, h, title=t, body=d, theme=theme,
        title_size=13, body_size=10.5, align=PP_ALIGN.LEFT)
    if i < n - 1:
        arrow(s, x + 0.55, y + h, x + 0.55, y + h + gap, width=1.4)
    y += h + gap
# re-plan loop callout on the right
box(s, 8.45, 2.95, 4.3, 1.55,
    title="Re-plan loop (autonomous)",
    body=["If the gate says “not done” it loops back to the Planner",
          "with the current blackboard snapshot.",
          "Guard: iteration ≤ max_replans (default 3).",
          "At the cap → forced Synthesizer with a partial-result flag."],
    theme=INDIGO, title_size=13, body_size=11, align=PP_ALIGN.LEFT)
box(s, 8.45, 4.7, 4.3, 1.75,
    title="Why a shared blackboard",
    body=["One workspace every agent reads and writes — no direct calls.",
          "• Decouples agents: a new one just reads/writes keys.",
          "• Drives the waves: an agent runs once its inputs appear.",
          "• Error flags make partial results visible to the gate.",
          "• Auditable + replayable — every result in one place."],
    theme=TEAL, title_size=13, body_size=10.5, align=PP_ALIGN.LEFT)
footer(s, 3)

# ============================================================================
# SLIDE 4 — How the Planner builds the plan
# ============================================================================
s = slide()
header(s, "How the Planner builds the plan",
       "Registry-driven — never a hardcoded list of agents")
nums = [
    ("1", "Split the prompt into intents",
     "Each clause of the request becomes a candidate subtask."),
    ("2", "Match each intent to a registered agent",
     "Compare the phrase to every agent’s capability description + tags."),
    ("3", "Wire data dependencies into a DAG",
     "If t-B needs t-A’s output, add an edge → ordering falls out."),
    ("4", "Assign a timeout budget per subtask",
     "Pull the SLA / timeout hint from the registry record."),
]
y = 1.85
for num, t, d in nums:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y),
                           Inches(0.62), Inches(0.62))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    c.line.fill.background(); c.shadow.inherit = False
    _set_text(c.text_frame, num, 20, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, 1.75, y + 0.02, 6.0, 0.45, t, size=16, color=INK, bold=True)
    textbox(s, 1.75, y + 0.5, 6.0, 0.5, d, size=12, color=INK2)
    y += 1.18
box(s, 8.2, 1.85, 4.55, 4.55,
    title="Re-entry on re-plan",
    body=["On a re-plan loop the Planner receives the current",
          "blackboard snapshot and revises only what is missing.",
          "",
          "Match is by meaning, not keywords:",
          "• “related outages” → Historical Comparison Agent",
          "   (not Linked Outage or OASIS Comparison).",
          "",
          "Version is pinned at plan time → reproducible replay.",
          "SLA hint flows to the Orchestrator for per-subtask",
          "budgets and transport choice.",
          "",
          "A new agent is usable the moment it registers —",
          "no Planner code change required."],
    theme=NEUTRAL, title_size=15, body_size=11.5, align=PP_ALIGN.LEFT)
footer(s, 4)

# ============================================================================
# SLIDE 5 — Worked example: DAG of an outage query
# ============================================================================
s = slide()
header(s, "Worked example — the plan as a DAG",
       "“For outage 17299126: details, related outages, scan ERC, OASIS derates, NRS-Scott loss”")
# Wave 1
chip(s, 0.9, 1.75, "WAVE 1", AMBER, w=1.05, size=10)
box(s, 0.9, 2.12, 5.6, 0.62, title="t1 · Database Agent",
    body=["“details” → pull outage 17299126  ·  deps: none  ·  SLA < 2 s"],
    theme=TERRA, title_size=13, body_size=10.5)
arrow(s, 3.7, 2.74, 3.7, 3.05, width=1.4)
# Wave 2 - three parallel
chip(s, 0.9, 3.05, "WAVE 2  —  parallel", AMBER, w=2.3, size=10)
w2 = [
    ("t2 · RAG Agent", "region + NRS-Scott", "deps: t1 · SLA < 4 s"),
    ("t3 · Historical", "find related outages", "deps: t1 · SLA < 5 s"),
    ("t4 · OASIS", "path derates", "deps: t1 · SLA < 8 s"),
]
bx = 0.9
for t, a, b in w2:
    box(s, bx, 3.42, 1.78, 0.95, title=t, body=[a, b], theme=TERRA,
        title_size=12, body_size=9.5)
    bx += 1.91
arrow(s, 3.7, 4.37, 3.7, 4.68, width=1.4)
textbox(s, 3.85, 4.4, 3, 0.3, "needs t3.related_list", size=10, color=INK2)
# Wave 3
chip(s, 0.9, 4.68, "WAVE 3  —  waits for t3", INDIGO, w=2.6, size=10)
box(s, 1.6, 5.05, 4.2, 0.62, title="t5 · Keyword Scanning Agent",
    body=["scan related outages for “ERC”  ·  deps: t3  ·  SLA < 6 s"],
    theme=TERRA, title_size=13, body_size=10)
# Right: how the match works
box(s, 7.1, 1.75, 5.65, 4.7,
    title="How the match works",
    body=["The Planner reads the registry — never a hardcoded list.",
          "For each intent it compares the phrase to every agent’s",
          "registered capability description + tags.",
          "",
          "• t2 and t4 each depend only on t1’s outage record, so",
          "   they run in parallel alongside t3 in Wave 2.",
          "• t5 reads t3’s related list → depends_on t3 → Wave 3.",
          "",
          "• The SLA hint is passed to the Planner so the",
          "   Orchestrator can enforce per-subtask budgets and",
          "   choose the transport.",
          "• Agent version is pinned at plan time → reproducible replay.",
          "• A new agent is usable as soon as it registers."],
    theme=NEUTRAL, title_size=15, body_size=11, align=PP_ALIGN.LEFT)
footer(s, 5)

# ============================================================================
# SLIDE 6 — Plan -> execution waves (generic) + orchestrator
# ============================================================================
s = slide()
header(s, "Plan → execution waves",
       "Parallelism falls out of the graph — the Planner never says “run in parallel”")
# generic DAG (left)
textbox(s, 0.9, 1.7, 5, 0.35, "Plan as a dependency graph", size=14,
        color=INK, bold=True)
box(s, 0.9, 2.15, 1.7, 0.55, title="Task A", theme=TERRA, title_size=14)
box(s, 3.0, 2.15, 1.7, 0.55, title="Task B", theme=TERRA, title_size=14)
chip(s, 0.9, 2.05, "", AMBER, w=0.01, h=0.01)  # spacer no-op
arrow(s, 1.75, 2.7, 1.75, 3.05, width=1.4)
arrow(s, 3.85, 2.7, 3.85, 3.05, width=1.4)
box(s, 0.9, 3.05, 1.7, 0.55, title="Task C", theme=TERRA, title_size=14)
box(s, 3.0, 3.05, 1.7, 0.55, title="Task D", theme=TERRA, title_size=14)
arrow(s, 1.75, 3.6, 2.62, 3.95, width=1.4)
arrow(s, 3.85, 3.6, 2.98, 3.95, width=1.4)
box(s, 1.95, 3.95, 1.7, 0.55, title="Task E", theme=TERRA, title_size=14)
textbox(s, 0.9, 4.7, 5.7, 1.6, [
    "A, B have no deps → Wave 1 (parallel).",
    "C needs A; D needs B → Wave 2 (parallel).",
    "E needs both C and D → Wave 3 (join, sequential wait)."],
    size=12, color=INK2, space_after=5)
# orchestrator execution (right)
textbox(s, 7.0, 1.7, 5, 0.35, "Orchestrator execution", size=14,
        color=INK, bold=True)
box(s, 7.0, 2.15, 5.7, 0.6, title="Wave 1 — parallel",
    body=["A and B run at once"], theme=TEAL, title_size=14, body_size=11)
arrow(s, 9.85, 2.75, 9.85, 3.05, width=1.4)
box(s, 7.0, 3.05, 5.7, 0.6, title="Wave 2 — parallel",
    body=["C and D run at once"], theme=TEAL, title_size=14, body_size=11)
arrow(s, 9.85, 3.65, 9.85, 3.95, width=1.4)
box(s, 7.0, 3.95, 5.7, 0.6, title="Wave 3 — sequential",
    body=["E waits for C and D"], theme=INDIGO, title_size=14, body_size=11)
box(s, 7.0, 4.8, 5.7, 1.5, title="Bounds / policy promotion",
    body=["A subtask whose estimated duration exceeds the budget",
          "(e.g. OASIS > 20 s) is automatically promoted from",
          "sync JSON-RPC to async Kafka by an Orchestrator rule."],
    theme=AMBER, title_size=13, body_size=11, align=PP_ALIGN.LEFT)
footer(s, 6)

# ============================================================================
# SLIDE 7 — Latency & timeout budget table
# ============================================================================
s = slide()
header(s, "Latency & timeout budget",
       "Per-subtask SLA hint drives transport choice and the on-timeout policy")
rows = [
    ("Agent", "SLA hint", "Transport", "On timeout"),
    ("t1  Database", "< 2 s", "JSON-RPC", "retry ×3 → fail"),
    ("t2  RAG", "< 4 s", "JSON-RPC", "retry ×3 → fail"),
    ("t3  Historical", "< 5 s", "JSON-RPC", "retry ×2 → fail"),
    ("t4  OASIS (slow)", "< 8 s", "Kafka async", "partial result"),
    ("t5  Keyword", "< 6 s", "JSON-RPC", "retry ×2 → fail"),
]
tx, ty = 0.9, 1.9
tw = 8.0
colw = [3.0, 1.5, 1.8, 1.7]
rh = 0.52
for r, row in enumerate(rows):
    cx = tx
    head = (r == 0)
    for c, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(ty + r*rh),
                                  Inches(colw[c]), Inches(rh))
        cell.fill.solid()
        cell.fill.fore_color.rgb = (RGBColor(0x53,0x4A,0xB7) if head
                                    else (WHITE if r % 2 else RGBColor(0xF4,0xF2,0xEC)))
        cell.line.color.rgb = RGBColor(0xDD,0xDB,0xD2); cell.line.width = Pt(0.75)
        cell.shadow.inherit = False
        _set_text(cell.text_frame, val, 12 if head else 11.5,
                  WHITE if head else INK, head,
                  PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        cell.text_frame.margin_left = Inches(0.12)
        cx += colw[c]
box(s, 9.3, 1.9, 3.45, 3.4,
    title="Why it matters",
    body=["• The SLA hint is a registry field, not a guess.",
          "",
          "• Blocking dependency → synchronous JSON-RPC.",
          "",
          "• Slow / fan-out work → asynchronous Kafka.",
          "",
          "• t4 exceeds the 20 s estimate, so the",
          "   Orchestrator promotes it to Kafka async",
          "   under a bounds/policy rule.",
          "",
          "• On timeout, each agent has an explicit",
          "   retry budget → then fail or partial."],
    theme=TEAL, title_size=14, body_size=11, align=PP_ALIGN.LEFT)
footer(s, 7)

# ============================================================================
# SLIDE 8 — Canonical patterns mapping
# ============================================================================
s = slide()
header(s, "Which canonical pattern each stage uses",
       "Each well-known agent pattern maps to one stage of the flow")
pairs = [
    ("Planning", "Planner emits the DAG"),
    ("ReAct", "Autonomous mode loop"),
    ("Tool use", "Tool / MCP gateway"),
    ("Reflection", "Completion gate + reflect"),
    ("Multi-agent", "Agent team + blackboard"),
]
y = 1.9
for pat, stage in pairs:
    box(s, 1.4, y, 3.0, 0.74, title=pat, theme=GREEN, title_size=16)
    arrow(s, 4.45, y + 0.37, 5.25, y + 0.37, width=1.6)
    box(s, 5.3, y, 5.4, 0.74, title=stage, theme=TEAL, title_size=16)
    y += 0.92
footer(s, 8)

# ============================================================================
# SLIDE 9 — A2A transport
# ============================================================================
s = slide()
header(s, "Agent-to-agent (A2A) transport",
       "Both supported — the Orchestrator derives the transport, it is not configured per agent")
box(s, 0.9, 1.75, 11.85, 0.7, title="Transport interface (platform)",
    body=["agent.call(target, payload)   ·   agent.publish(topic, event)   ·   agent.subscribe(topic)"],
    theme=INDIGO, title_size=14, body_size=12)
# Two transport cards
box(s, 0.9, 2.75, 5.75, 3.55,
    title="HTTPS / JSON-RPC",
    body=["Synchronous · point-to-point · low latency",
          "",
          "Use when one agent needs a specific answer now",
          "and will wait for it (a blocking call).",
          "",
          "• Orchestrator → Database Agent for t1.outage",
          "• Correlation id, request → response, easy to trace",
          "• Natural fit for dependency waves that block"],
    theme=TEAL, title_size=16, body_size=12, align=PP_ALIGN.LEFT)
box(s, 7.0, 2.75, 5.75, 3.55,
    title="Kafka",
    body=["Asynchronous · pub/sub · durable + replayable",
          "",
          "Use for one-to-many, fire-and-forget, or long jobs.",
          "",
          "• Monitor Agent emits events many consumers tail",
          "• Long OASIS scan publishes ‘done’ vs holding HTTP",
          "• Outage-change events fan out to interested agents",
          "• Backpressure: check consumer lag before publishing"],
    theme=AMBER, title_size=16, body_size=12, align=PP_ALIGN.LEFT)
textbox(s, 0.9, 6.45, 11.8, 0.4,
        [("Rule of thumb:  ", INK, True, 12.5),
         ("need a specific answer and will wait → JSON-RPC   ·   "
          "one-to-many / fire-and-forget / durable → Kafka", INK2, False, 12.5)],
        size=12.5)
footer(s, 9)

# ============================================================================
# SLIDE 10 — Memory & state
# ============================================================================
s = slide()
header(s, "Memory & state",
       "Private scratchpad vs a governed, shared platform")
box(s, 0.9, 1.75, 11.85, 0.72,
    title="Per-agent scratchpad (private, in-step)",
    body=["ReAct reasoning and drafts — never a source of truth. Only the final schema-valid output is promoted."],
    theme=NEUTRAL, title_size=14, body_size=11.5)
box(s, 0.9, 2.65, 11.85, 0.82,
    title="Memory / State service (governed interface)",
    body=["state.put / state.get     ·     memory.write / memory.query",
          "enforces per-run + per-app scoping, routing, and audit"],
    theme=INDIGO, title_size=14, body_size=11.5)
stores = [
    ("Redis", "Working state", ["blackboard, per-run,", "hot, TTL / expire", "", "per-run · shared in run"], TEAL),
    ("Postgres", "Relational long-term", ["past analyses, entity", "links, audit trail", "", "cross-run · shared"], TEAL),
    ("Milvus", "Semantic long-term", ["embeddings: related /", "similar outages, RAG", "", "cross-run · shared"], TEAL),
]
bx = 0.9
for name, sub, body, theme in stores:
    box(s, bx, 3.7, 3.8, 1.95, title=name,
        body=[sub] + body, theme=theme, title_size=16, body_size=11,
        align=PP_ALIGN.LEFT)
    bx += 3.97
box(s, 0.9, 5.85, 11.85, 1.0,
    title="Write discipline",
    body=["During a run, agents write working state (Redis) freely. Durable writes (Postgres / Milvus) commit",
          "only after the completion gate passes — a failed or re-planned run leaves no half-baked facts."],
    theme=AMBER, title_size=14, body_size=11.5)
footer(s, 10)

# ============================================================================
# SLIDE 11 — Commit step
# ============================================================================
s = slide()
header(s, "Commit step — what the Synthesizer persists",
       "The schema decides routing, not the LLM — the Synthesizer only filters, dedupes, and commits")
flow = [
    ("Blackboard outputs", "details, related, hits, derates", TEAL),
    ("Keep only persist = true", "set in output_schema at design time", NEUTRAL),
    ("Dedupe / upsert by key", "merge against existing records", NEUTRAL),
    ("memory.write(namespace, record)", "one record → fans out to both stores", INDIGO),
]
y = 1.85
for t, d, theme in flow:
    box(s, 1.2, y, 5.4, 0.74, title=t, body=[d], theme=theme,
        title_size=14, body_size=11, align=PP_ALIGN.LEFT)
    if y < 4.0:
        arrow(s, 1.9, y + 0.74, 1.9, y + 0.92, width=1.4)
    y += 0.92
arrow(s, 6.6, 4.55, 7.7, 4.0, width=1.4)
arrow(s, 6.6, 4.55, 7.7, 5.1, width=1.4)
box(s, 7.7, 3.6, 5.0, 0.85, title="Postgres",
    body=["structured fields → a row (ids, link type, timestamps)"],
    theme=TEAL, title_size=14, body_size=11, align=PP_ALIGN.LEFT)
box(s, 7.7, 4.7, 5.0, 0.85, title="Milvus",
    body=["text field → embedding vector (shared key links it to the row)"],
    theme=TEAL, title_size=14, body_size=11, align=PP_ALIGN.LEFT)
textbox(s, 1.2, 5.45, 5.5, 1.4, [
    ("In the example:", INK, True, 12.5),
    ("• New outage link → persist → Postgres row + Milvus embedding.", INK2, False, 11.5),
    ("• Details, ERC hits, derates → ephemeral, stay in Redis.", INK2, False, 11.5),
    ("• Generated concerns → persist only if you cache analyses.", INK2, False, 11.5),
], size=12, space_after=4)
footer(s, 11)

# ============================================================================
# SLIDE 12 — Completion gate & synthesizer
# ============================================================================
s = slide()
header(s, "Completion gate & Synthesizer",
       "The decision point, then the single composed answer")
box(s, 0.9, 1.8, 5.75, 4.4,
    title="What the completion gate decides",
    body=["Inspects the blackboard and answers one question: done?",
          "",
          "• Bounds first — stop if steps / cost / time exceeded.",
          "• Plan satisfied — every subtask’s output_schema met.",
          "• Goal met — reflection check on the combined result.",
          "• Re-plan guard — iteration counter checked before loop.",
          "",
          "   iter < max  → re-plan with the current snapshot.",
          "   iter = max → forced Synthesizer, partial-result flag."],
    theme=INDIGO, title_size=15, body_size=12, align=PP_ALIGN.LEFT)
box(s, 7.0, 1.8, 5.75, 4.4,
    title="What the Synthesizer produces",
    body=["Reads the whole blackboard and composes ONE answer.",
          "",
          "• Merges details, related, keyword hits, derates, concerns.",
          "• Resolves overlaps and orders the response for the user.",
          "• Marks sections [PARTIAL] where agents errored out.",
          "",
          "Does not invent routing — the output_schema decides what",
          "is persisted; the Synthesizer filters, dedupes, commits."],
    theme=TEAL, title_size=15, body_size=12, align=PP_ALIGN.LEFT)
footer(s, 12)

# ============================================================================
# SLIDE 13 — Observability
# ============================================================================
s = slide()
header(s, "Observability — tracing, metrics & audit",
       "Every component emits spans, collected into a unified trace per run_id")
box(s, 0.9, 1.8, 11.85, 0.72, title="Trace collector (platform-wide)",
    body=["run_id → parent span → child spans per agent / tool call  ·  OpenTelemetry-compatible"],
    theme=INDIGO, title_size=14, body_size=12)
spans = [
    ("Planner span", ["intent count,", "agents selected,", "DAG depth", "~200 ms"], NEUTRAL),
    ("Agent span (×N)", ["agent_id, version,", "wave, status,", "retry_count", "per SLA hint"], TERRA),
    ("Tool span (×M)", ["tool_name, target,", "status, tokens_used,", "input/output hash", "API-dependent"], AMBER),
    ("Gate + Synth span", ["iter_count, goal_met,", "partial_flags,", "tokens_committed", "~300 ms"], TEAL),
]
bx = 0.9
for t, body, theme in spans:
    box(s, bx, 2.85, 2.85, 1.9, title=t, body=body, theme=theme,
        title_size=13, body_size=10.5, align=PP_ALIGN.LEFT)
    bx += 2.98
box(s, 0.9, 5.0, 11.85, 0.72, title="Downstream sinks",
    body=["Spans flow to MLflow (and any OpenTelemetry-compatible backend) for dashboards, diffing, and audit."],
    theme=NEUTRAL, title_size=14, body_size=11.5)
footer(s, 13)

# ============================================================================
# SLIDE 14 — Agent versioning
# ============================================================================
s = slide()
header(s, "Agent versioning — the registry contract",
       "Version is pinned at plan time, stored in each span, and written to the Postgres audit row")
rows = [
    ("Field", "Example", "Notes"),
    ("agent_id", "historical-comparison", "stable identifier"),
    ("version", "2.3.1", "semver, pinned at plan time"),
    ("capability_tags", "[“related_outages”, “rag”]", "matched by the Planner"),
    ("input_schema", "{ outage_id, region }", "validated before dispatch"),
    ("output_schema", "{ related: [ ] }", "persist flag per field"),
    ("sla_ms", "5000", "Orchestrator timeout budget"),
    ("transport", "json-rpc", "or kafka, or both"),
    ("status", "active | deprecated", "Planner skips deprecated"),
]
tx, ty = 0.9, 1.85
colw = [2.6, 3.4, 3.6]
rh = 0.44
for r, row in enumerate(rows):
    cx = tx
    head = (r == 0)
    for c, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(ty + r*rh),
                                  Inches(colw[c]), Inches(rh))
        cell.fill.solid()
        cell.fill.fore_color.rgb = (RGBColor(0x85,0x4F,0x0B) if head
                                    else (WHITE if r % 2 else RGBColor(0xF7,0xF1,0xE6)))
        cell.line.color.rgb = RGBColor(0xDD,0xDB,0xD2); cell.line.width = Pt(0.75)
        cell.shadow.inherit = False
        _set_text(cell.text_frame, val, 11.5, WHITE if head else INK, head,
                  PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        cell.text_frame.margin_left = Inches(0.12)
        cx += colw[c]
textbox(s, 0.9, 5.95, 11.8, 1.1, [
    ("• New version → register as a new record; the old one stays until deprecated.", INK2, False, 11.5),
    ("• Planner pins version at plan time → a run is reproducible by replaying the same plan JSON.", INK2, False, 11.5),
    ("• Breaking output_schema change → major bump; the Planner will not auto-select across majors.", INK2, False, 11.5),
    ("• Trace diff: compare two run_ids to see which agent version changed and how output differed.", INK2, False, 11.5),
], size=11.5, space_after=3)
footer(s, 14)

# ============================================================================
# SLIDE 15 — Partial streaming & Kafka backpressure
# ============================================================================
s = slide()
header(s, "Partial streaming & Kafka backpressure",
       "Progressive delivery to the user, with consumer-lag control on the bus")
box(s, 0.9, 1.85, 5.85, 4.5,
    title="Partial streaming",
    body=["The Synthesizer does not wait for all waves. As each wave",
          "completes and writes to the blackboard, a wave-done event",
          "is published; the Synthesizer flushes a partial section.",
          "",
          "• Wave 1 done → stream outage details immediately.",
          "• Wave 2 done → append related outages + derates.",
          "• Wave 3 done → append ERC keyword hits. Final flush.",
          "",
          "• Wave errors → stream a [PARTIAL] placeholder at once;",
          "   the user sees available data without waiting for the",
          "   full retry cycle."],
    theme=TEAL, title_size=16, body_size=11.5, align=PP_ALIGN.LEFT)
box(s, 7.0, 1.85, 5.75, 4.5,
    title="Kafka backpressure",
    body=["Problem: if producers outpace consumers, lag grows",
          "unboundedly — delaying agents, bloating the broker,",
          "risking OOM.",
          "",
          "Controls:",
          "• Lag alert when consumer lag crosses a threshold.",
          "• Producer pause: the Orchestrator checks lag before",
          "   publishing the next wave → exponential backoff.",
          "• Dead-letter queue on max retry.",
          "• Replayability: Kafka retains messages — a recovered",
          "   consumer replays from the last committed offset."],
    theme=AMBER, title_size=16, body_size=11.5, align=PP_ALIGN.LEFT)
footer(s, 15)

# ============================================================================
# SLIDE 16 — Closing
# ============================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.1), SW, Inches(3.3))
band.fill.solid(); band.fill.fore_color.rgb = WHITE
band.line.fill.background(); band.shadow.inherit = False
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.35),
                         Inches(0.12), Inches(2.8))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background(); bar.shadow.inherit = False
textbox(s, 1.35, 2.3, 11, 0.8, "The shape of the platform", size=32,
        color=INK, bold=True)
textbox(s, 1.38, 3.15, 11.2, 2.0, [
    ("•  The Planner reads a registry and emits a DAG — parallelism is implied by data dependencies.", INK2, False, 14),
    ("•  The Orchestrator runs the DAG in waves and derives the A2A transport per subtask.", INK2, False, 14),
    ("•  A shared blackboard decouples agents; durable writes commit only after the gate passes.", INK2, False, 14),
    ("•  Spans, versioning, streaming, and backpressure make runs observable, reproducible, and resilient.", INK2, False, 14),
], size=14, space_after=10, line_spacing=1.05)
footer(s, 16)

out = r"d:\BaseAgentFramework\docs\Agentic_Framework_Flow_Plan_Patterns.pptx"
prs.save(out)
print("Saved:", out, "with", len(prs.slides._sldIdLst), "slides")
